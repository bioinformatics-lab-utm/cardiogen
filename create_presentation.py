#!/usr/bin/env python3
"""
Create PowerPoint Presentation from SRA Metadata Analysis
Automatically generates a comprehensive presentation with all visualizations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
from datetime import datetime
import os


class SRAPresentationCreator:
    """Create PowerPoint presentation from analysis figures."""
    
    def __init__(self, figures_dir: str, output_file: str = None):
        """
        Initialize presentation creator.
        
        Args:
            figures_dir: Directory containing figure PNG files
            output_file: Output PowerPoint file path
        """
        self.figures_dir = Path(figures_dir)
        self.output_file = output_file or f'SRA_Metadata_Analysis_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        
        # Define slide sections
        self.sections = {
            'intro': {
                'title': 'Introduction',
                'figures': []
            },
            'instruments': {
                'title': 'Sequencing Instruments',
                'figures': ['01_instrument_distribution_bar.png', 
                           '02_instrument_distribution_pie.png']
            },
            'library': {
                'title': 'Library Preparation Strategies',
                'figures': ['03_library_strategy_bar.png',
                           '04_library_strategy_pie.png']
            },
            'depth': {
                'title': 'Sequencing Depth Analysis',
                'figures': ['05_sequencing_depth_histogram.png',
                           '05b_sequencing_depth_violin.png',
                           '05c_total_bases_histogram.png']
            },
            'demographics': {
                'title': 'Sample Demographics',
                'figures': ['06_sex_distribution.png',
                           '06b_sex_distribution_pie.png']
            },
            'completeness': {
                'title': 'Metadata Completeness',
                'figures': ['07_metadata_completeness.png',
                           '23_completeness_distribution.png',
                           '24_completeness_categories.png']
            },
            'organizations': {
                'title': 'Contributing Organizations',
                'figures': ['08_top_organizations.png']
            },
            'correlations': {
                'title': 'Data Correlations',
                'figures': ['09_correlation_matrix.png']
            },
            'diseases': {
                'title': 'Disease Distribution',
                'figures': ['10_disease_distribution_bar.png',
                           '11_disease_distribution_pie.png']
            },
            'tissues': {
                'title': 'Tissue Types',
                'figures': ['12_tissue_distribution_bar.png',
                           '13_tissue_distribution_pie.png']
            },
            'scatter': {
                'title': 'Correlation Analysis',
                'figures': ['17_scatter_spots_vs_bases.png']
            },
            'studies': {
                'title': 'Study-Level Analysis',
                'figures': ['18_study_size_distribution.png',
                           '19_top_studies.png']
            },
            'layout': {
                'title': 'Library Layout Comparison',
                'figures': ['20_library_layout_violin.png']
            },
            'patient_demographics': {
                'title': 'Patient Demographics Details',
                'figures': ['22_patient_demographics_completeness.png']
            },
            'columns': {
                'title': 'Column Categories',
                'figures': ['25_column_categories_bar.png',
                           '26_column_categories_pie.png']
            },
            'technical': {
                'title': 'Technical Parameters',
                'figures': ['27_library_layout.png',
                           '28_library_selection.png',
                           '29_library_source.png',
                           '30_platform_type_pie.png']
            },
            'quality': {
                'title': 'Quality Metrics',
                'figures': ['32_core_metadata_completeness.png']
            },
            'statistical': {
                'title': 'Statistical Summary',
                'figures': ['33_statistical_summary_heatmap.png']
            },
            'summary': {
                'title': 'Dataset Summary',
                'figures': ['34_dataset_summary.png']
            },
            'sample_level': {
                'title': 'Sample-Level Analysis (Patient Deduplication)',
                'figures': ['43_runs_per_sample_distribution.png',
                           '44_runs_vs_samples_comparison.png',
                           '45_top_samples_by_runs.png']
            },
            'phenotypes': {
                'title': 'Comprehensive Phenotype Analysis',
                'figures': ['46_phenotype_data_completeness.png',
                           '47_sex_distribution_all.png',
                           '48_age_distribution_histogram.png',
                           '49_disease_distribution_all.png',
                           '50_tissue_distribution_all.png',
                           '51_treatment_distribution_all.png',
                           '52_condition_distribution_all.png',
                           '53_diagnosis_distribution_all.png',
                           '54_disease_state_distribution_all.png',
                           '55_phenotype_diversity_all.png']
            }
        }
    
    def add_title_slide(self):
        """Add title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(20, 40, 80)
        
        # Add title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(14)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "SRA Metadata Analysis"
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.8), width, Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Comprehensive Analysis of Illumina Sequencing Data"
        
        p2 = subtitle_frame.paragraphs[0]
        p2.font.size = Pt(32)
        p2.font.color.rgb = RGBColor(200, 220, 255)
        p2.alignment = PP_ALIGN.CENTER
        
        # Add date
        date_box = slide.shapes.add_textbox(left, top + Inches(3.2), width, Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%B %d, %Y")
        
        p3 = date_frame.paragraphs[0]
        p3.font.size = Pt(20)
        p3.font.color.rgb = RGBColor(180, 200, 240)
        p3.alignment = PP_ALIGN.CENTER
        
        print("✓ Added title slide")
    
    def add_overview_slide(self):
        """Add overview slide with key statistics."""
        slide_layout = self.prs.slide_layouts[5]  # Title only
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Dataset Overview"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add statistics
        left = Inches(2)
        top = Inches(2)
        width = Inches(12)
        height = Inches(4.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        
        stats = [
            "📊 Total SRA Runs: 2,880",
            "👥 Unique Samples (Patients): 2,712",
            "📈 Average Runs per Sample: 1.06",
            "🔄 Samples with Multiple Runs: 87 (3.2%)",
            "📋 Total Metadata Columns: 300",
            "🧬 Phenotype Columns Analyzed: 27",
            "🔬 Platform: Illumina Only",
            "📊 Visualizations Generated: 45"
        ]
        
        for i, stat in enumerate(stats):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = stat
            p.font.size = Pt(28)
            p.font.bold = True
            p.space_before = Pt(12)
            p.font.color.rgb = RGBColor(40, 60, 100)
        
        print("✓ Added overview slide")
    
    def add_section_divider(self, section_title: str):
        """Add section divider slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(60, 80, 120)
        
        # Title
        left = Inches(1)
        top = Inches(3.5)
        width = Inches(14)
        height = Inches(2)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.text = section_title
        
        p = tf.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def add_figure_slide(self, figure_path: Path, title: str = None):
        """Add slide with a figure."""
        if not figure_path.exists():
            print(f"⚠️  Warning: {figure_path.name} not found, skipping...")
            return
        
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title if provided
        if title:
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(15)
            height = Inches(0.7)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title
            
            p = title_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(40, 60, 100)
            p.alignment = PP_ALIGN.CENTER
        
        # Add figure - maximize size while maintaining aspect ratio
        img_left = Inches(0.5)
        img_top = Inches(1.2) if title else Inches(0.5)
        img_width = Inches(15)
        img_height = Inches(7.3) if title else Inches(8)
        
        try:
            pic = slide.shapes.add_picture(
                str(figure_path),
                img_left, img_top,
                width=img_width,
                height=img_height
            )
            print(f"  ✓ Added: {figure_path.name}")
        except Exception as e:
            print(f"  ✗ Error adding {figure_path.name}: {e}")
    
    def add_two_figures_slide(self, figure_paths: list, title: str = None):
        """Add slide with two figures side by side."""
        available_paths = [p for p in figure_paths if p.exists()]
        
        if not available_paths:
            print(f"⚠️  Warning: No figures available for {title}, skipping...")
            return
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title if provided
        if title:
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(15)
            height = Inches(0.7)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title
            
            p = title_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(40, 60, 100)
            p.alignment = PP_ALIGN.CENTER
        
        # Add figures
        img_top = Inches(1.2) if title else Inches(0.5)
        img_width = Inches(7.2)
        img_height = Inches(7.3) if title else Inches(8)
        
        for i, fig_path in enumerate(available_paths[:2]):
            img_left = Inches(0.5) if i == 0 else Inches(8.3)
            try:
                slide.shapes.add_picture(
                    str(fig_path),
                    img_left, img_top,
                    width=img_width,
                    height=img_height
                )
                print(f"  ✓ Added: {fig_path.name}")
            except Exception as e:
                print(f"  ✗ Error adding {fig_path.name}: {e}")
    
    def create_presentation(self):
        """Create complete presentation."""
        print("\n" + "="*60)
        print("Creating PowerPoint Presentation")
        print("="*60)
        
        # Title and overview
        self.add_title_slide()
        self.add_overview_slide()
        
        # Process each section
        for section_key, section_info in self.sections.items():
            if section_key == 'intro':
                continue
            
            section_title = section_info['title']
            figures = section_info['figures']
            
            if not figures:
                continue
            
            # Add section divider
            self.add_section_divider(section_title)
            print(f"\n📂 Section: {section_title}")
            
            # Add figures
            i = 0
            while i < len(figures):
                fig_paths = [self.figures_dir / figures[i]]
                
                # Check if we can fit two figures
                if i + 1 < len(figures):
                    next_fig = self.figures_dir / figures[i + 1]
                    # Add two figures if both exist
                    if next_fig.exists() and fig_paths[0].exists():
                        fig_paths.append(next_fig)
                        self.add_two_figures_slide(fig_paths, section_title)
                        i += 2
                        continue
                
                # Add single figure
                if fig_paths[0].exists():
                    self.add_figure_slide(fig_paths[0], section_title)
                i += 1
        
        # Save presentation
        self.prs.save(self.output_file)
        
        print("\n" + "="*60)
        print("✅ PRESENTATION CREATED!")
        print(f"📊 File: {self.output_file}")
        print(f"📁 Size: {os.path.getsize(self.output_file) / 1024 / 1024:.2f} MB")
        print(f"🎯 Slides: {len(self.prs.slides)}")
        print("="*60 + "\n")


def main():
    """Main entry point."""
    import argparse
    
    parser = argparse.ArgumentParser(
        description='Create PowerPoint presentation from SRA analysis figures'
    )
    
    parser.add_argument(
        '--figures', '-f',
        default='data/analysis_results/figures',
        help='Directory containing figure PNG files'
    )
    
    parser.add_argument(
        '--output', '-o',
        default=None,
        help='Output PowerPoint file path'
    )
    
    args = parser.parse_args()
    
    # Create presentation
    creator = SRAPresentationCreator(args.figures, args.output)
    creator.create_presentation()


if __name__ == '__main__':
    main()
